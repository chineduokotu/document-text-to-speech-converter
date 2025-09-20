"""
Document Readers Module

This module provides readers for various document formats including
PDF, DOCX, PPTX, TXT, and web content.
"""

import PyPDF2
from docx import Document
from pptx import Presentation
import chardet
from pathlib import Path
from typing import List, Optional, Union
import logging
import requests
from bs4 import BeautifulSoup

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class DocumentReader:
    """Base class for document readers."""
    
    @staticmethod
    def detect_encoding(file_path: Path) -> str:
        """Detect the encoding of a text file."""
        try:
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                return result['encoding'] or 'utf-8'
        except Exception as e:
            logger.warning(f"Could not detect encoding for {file_path}: {e}")
            return 'utf-8'


class TextFileReader(DocumentReader):
    """Reader for plain text files."""
    
    @staticmethod
    def read(file_path: Union[str, Path]) -> Optional[str]:
        """
        Read text from a plain text file.
        
        Args:
            file_path: Path to the text file
            
        Returns:
            str: The text content, or None if failed
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                logger.error(f"File not found: {file_path}")
                return None
            
            encoding = TextFileReader.detect_encoding(file_path)
            
            with open(file_path, 'r', encoding=encoding) as file:
                content = file.read()
            
            logger.info(f"Successfully read text file: {file_path}")
            return content.strip()
            
        except Exception as e:
            logger.error(f"Failed to read text file {file_path}: {e}")
            return None


class PDFReader(DocumentReader):
    """Reader for PDF files."""
    
    @staticmethod
    def read(file_path: Union[str, Path]) -> Optional[str]:
        """
        Read text from a PDF file.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            str: The extracted text content, or None if failed
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                logger.error(f"File not found: {file_path}")
                return None
            
            text_content = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                logger.info(f"PDF has {len(pdf_reader.pages)} pages")
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(f"--- Page {page_num + 1} ---\n{page_text}")
                        logger.debug(f"Extracted text from page {page_num + 1}")
                    except Exception as e:
                        logger.warning(f"Could not extract text from page {page_num + 1}: {e}")
            
            if not text_content:
                logger.warning(f"No text content extracted from PDF: {file_path}")
                return None
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Successfully read PDF file: {file_path}")
            return full_text.strip()
            
        except Exception as e:
            logger.error(f"Failed to read PDF file {file_path}: {e}")
            return None


class DOCXReader(DocumentReader):
    """Reader for Microsoft Word DOCX files."""
    
    @staticmethod
    def read(file_path: Union[str, Path]) -> Optional[str]:
        """
        Read text from a DOCX file.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            str: The extracted text content, or None if failed
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                logger.error(f"File not found: {file_path}")
                return None
            
            doc = Document(file_path)
            text_content = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_content.append("--- Table ---\n" + "\n".join(table_text))
            
            if not text_content:
                logger.warning(f"No text content extracted from DOCX: {file_path}")
                return None
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Successfully read DOCX file: {file_path}")
            return full_text.strip()
            
        except Exception as e:
            logger.error(f"Failed to read DOCX file {file_path}: {e}")
            return None


class PPTXReader(DocumentReader):
    """Reader for Microsoft PowerPoint PPTX files."""
    
    @staticmethod
    def read(file_path: Union[str, Path]) -> Optional[str]:
        """
        Read text from a PPTX file.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            str: The extracted text content, or None if failed
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                logger.error(f"File not found: {file_path}")
                return None
            
            presentation = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num + 1} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            if not text_content:
                logger.warning(f"No text content extracted from PPTX: {file_path}")
                return None
            
            full_text = "\n\n".join(text_content)
            logger.info(f"Successfully read PPTX file: {file_path}")
            return full_text.strip()
            
        except Exception as e:
            logger.error(f"Failed to read PPTX file {file_path}: {e}")
            return None


class WebContentReader(DocumentReader):
    """Reader for web content from URLs."""
    
    @staticmethod
    def read_url(url: str, timeout: int = 10) -> Optional[str]:
        """
        Read text content from a web URL.
        
        Args:
            url: The URL to read from
            timeout: Request timeout in seconds
            
        Returns:
            str: The extracted text content, or None if failed
        """
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            response = requests.get(url, headers=headers, timeout=timeout)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text content
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            if not text.strip():
                logger.warning(f"No text content extracted from URL: {url}")
                return None
            
            logger.info(f"Successfully read web content from: {url}")
            return text.strip()
            
        except Exception as e:
            logger.error(f"Failed to read web content from {url}: {e}")
            return None


class UniversalDocumentReader:
    """Universal document reader that automatically detects file type and uses appropriate reader."""
    
    SUPPORTED_EXTENSIONS = {
        '.txt': TextFileReader,
        '.pdf': PDFReader,
        '.docx': DOCXReader,
        '.pptx': PPTXReader,
    }
    
    @classmethod
    def read_file(cls, file_path: Union[str, Path]) -> Optional[str]:
        """
        Read content from any supported file type.
        
        Args:
            file_path: Path to the file
            
        Returns:
            str: The extracted text content, or None if failed
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                logger.error(f"File not found: {file_path}")
                return None
            
            extension = file_path.suffix.lower()
            
            if extension not in cls.SUPPORTED_EXTENSIONS:
                logger.error(f"Unsupported file type: {extension}")
                logger.info(f"Supported types: {list(cls.SUPPORTED_EXTENSIONS.keys())}")
                return None
            
            reader_class = cls.SUPPORTED_EXTENSIONS[extension]
            return reader_class.read(file_path)
            
        except Exception as e:
            logger.error(f"Failed to read file {file_path}: {e}")
            return None
    
    @classmethod
    def read_url(cls, url: str) -> Optional[str]:
        """
        Read content from a web URL.
        
        Args:
            url: The URL to read from
            
        Returns:
            str: The extracted text content, or None if failed
        """
        return WebContentReader.read_url(url)
    
    @classmethod
    def get_supported_extensions(cls) -> List[str]:
        """Get list of supported file extensions."""
        return list(cls.SUPPORTED_EXTENSIONS.keys())
    
    @classmethod
    def is_supported(cls, file_path: Union[str, Path]) -> bool:
        """Check if a file type is supported."""
        file_path = Path(file_path)
        return file_path.suffix.lower() in cls.SUPPORTED_EXTENSIONS


def main():
    """Test the document readers."""
    print("Document Readers Test")
    print("====================")
    
    # Test different file types if they exist
    test_files = [
        "sample.txt",
        "sample.pdf", 
        "sample.docx",
        "sample.pptx"
    ]
    
    reader = UniversalDocumentReader()
    
    print(f"Supported extensions: {reader.get_supported_extensions()}")
    
    for file_name in test_files:
        if Path(file_name).exists():
            print(f"\nTesting {file_name}:")
            content = reader.read_file(file_name)
            if content:
                print(f"Content preview: {content[:100]}...")
            else:
                print("Failed to read content")
        else:
            print(f"\nSkipping {file_name} (file not found)")
    
    # Test URL reading
    test_url = "https://httpbin.org/html"
    print(f"\nTesting URL: {test_url}")
    try:
        content = reader.read_url(test_url)
        if content:
            print(f"URL content preview: {content[:100]}...")
        else:
            print("Failed to read URL content")
    except Exception as e:
        print(f"URL test failed: {e}")


if __name__ == "__main__":
    main()